#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
更新比赛 PPT（在原文件末尾追加幻灯片）。
原版已备份为：2_PPT_*_原版备份_20260525.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

PPT_PATH = Path("/srv/supercare/比赛文档/2_PPT_超级[GP-护士-照护员]智能体协同算法.pptx")
NARRATIVE_PNG = Path("/srv/supercare/task-agent/output/整体叙事图_本地序贯分期与超级GP协同.png")


def _add_title_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layouts = list(prs.slide_layouts)
    layout = layouts[1] if len(layouts) > 1 else layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is None and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
    if body:
        tf = body.text_frame
        tf.clear()
        for index, line in enumerate(bullets):
            p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)


def _pick_blank_layout(prs: Presentation):
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise RuntimeError("PPT 无可用版式")
    for layout in layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "空白" in name:
            return layout
    return layouts[0]


def _add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(_pick_blank_layout(prs))
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.8))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    tx.text_frame.paragraphs[0].font.bold = True
    if image_path.is_file():
        slide.shapes.add_picture(str(image_path), Inches(0.35), Inches(1.0), width=Inches(9.3))


def main() -> None:
    if not PPT_PATH.is_file():
        raise FileNotFoundError(PPT_PATH)
    prs = Presentation(str(PPT_PATH))
    before = len(prs.slides)

    _add_title_slide(
        prs,
        "创新：A0 本地序贯分期 + 贝叶斯",
        [
            "DataSource 10 例 → 生物标志物 Excel（本地解析）",
            "本地序贯疾病分期：阶段 1–5、亚型 1–3、序贯事件",
            "贝叶斯后验：BPSD 升级 / 跌倒 / 照护强度 + 决策阈值",
            "工具注入 A3 超级 GP、A6 任务包、A8 升级、A15 分期报告",
            "全程本地运行，无外部第三方分期 API",
        ],
    )
    _add_title_slide(
        prs,
        "超级 GP 针对性方案（示例：陈女士）",
        [
            "阶段 2 轻度进展期 · 亚型 1 认知-行为主导",
            "BPSD 30 天升级后验约 5%（低于会诊阈值）",
            "产出：a3_GP协作专业答复.pdf",
            "产出：a15_老年健康序贯分期与贝叶斯风险报告.pdf",
            "命令：python task-agent/run_staging_pipeline.py",
        ],
    )
    if NARRATIVE_PNG.is_file():
        _add_picture_slide(prs, "整体叙事图（本地序贯分期 → 贝叶斯 → 超级GP）", NARRATIVE_PNG)

    prs.save(str(PPT_PATH))
    print(f"PPT 已更新: {PPT_PATH}")
    print(f"幻灯片: {before} → {len(prs.slides)}（新增 {len(prs.slides) - before} 页）")
    print("原版备份: 2_PPT_*_原版备份_20260525.pptx")


if __name__ == "__main__":
    main()
