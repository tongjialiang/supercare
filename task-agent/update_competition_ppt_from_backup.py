#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
从原版备份 PPT 另存为「含序贯贝叶斯」版：在摘要/创新/技术方案/场景等页增补内容，
并在技术方案后插入算法、公式、先验文献、配图幻灯片。不覆盖原版备份。
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from staging.prior_literature import (
    PRIOR_ASSUMPTION_H,
    PRIOR_FORMULA_NOTE,
    PRIOR_REFERENCES,
    STAGE_BPSD_PRIORS_LITERATURE,
    STAGE_LABELS,
    STAGE_PRIOR_DERIVATION,
)

COMPETITION = Path("/srv/supercare/比赛文档")
BACKUP_PPT = COMPETITION / "2_PPT_超级[GP-护士-照护员]智能体协同算法_原版备份_20260525.pptx"
OUTPUT_PPT = COMPETITION / "2_PPT_超级[GP-护士-照护员]智能体协同算法_含序贯贝叶斯_20260525.pptx"
TASK_OUTPUT = Path(__file__).resolve().parent / "output"

FIGURES = [
    ("A0 整体叙事：DataSource → 序贯分期 → 贝叶斯 → 超级GP", "整体叙事图_本地序贯分期与超级GP协同.png"),
    ("疾病进展曲线（横轴=共同分期，纵轴=归一化进展指数）", "a0_疾病进展曲线_陈女士.png"),
    ("序贯事件槽位与 MCMC 整体分期", "a0_序贯事件进展图_陈女士.png"),
    ("贝叶斯先验 vs 后验（文献先验 + 标志物似然）", "a0_贝叶斯先验后验图_陈女士.png"),
]


def _pick_title_body_layout(prs: Presentation):
    layouts = list(prs.slide_layouts)
    return layouts[1] if len(layouts) > 1 else layouts[0]


def _pick_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "空白" in name:
            return layout
    return prs.slide_layouts[0]


def move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    """将幻灯片从 old_index 移动到 new_index（0-based）。"""
    slide_id_list = presentation.slides._sldIdLst
    elements = list(slide_id_list)
    target = elements[old_index]
    slide_id_list.remove(target)
    slide_id_list.insert(new_index, target)


def _find_thanks_slide_index(prs: Presentation) -> int | None:
    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "谢谢" in (shape.text_frame.text or ""):
                return index
    return None


def _append_paragraph(text_frame, text: str, font_size: int = 14, bold: bool = False) -> None:
    paragraph = text_frame.add_paragraph()
    paragraph.text = text
    paragraph.level = 0
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold


def _add_textbox_slide(
    prs: Presentation,
    title: str,
    lines: list[str],
    font_size: int = 16,
) -> None:
    layout = _pick_title_body_layout(prs)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is None and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
    if body:
        text_frame = body.text_frame
        text_frame.clear()
        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = Pt(font_size)
    else:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(font_size)


def _add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(_pick_blank_layout(prs))
    textbox = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(9.3), Inches(0.75))
    textbox.text_frame.text = title
    textbox.text_frame.paragraphs[0].font.size = Pt(22)
    textbox.text_frame.paragraphs[0].font.bold = True
    if image_path.is_file():
        slide.shapes.add_picture(str(image_path), Inches(0.3), Inches(0.95), width=Inches(9.4))


def _add_callout_to_slide(slide, top_inch: float, lines: list[str]) -> None:
    """在已有幻灯片底部增加说明框。"""
    box = slide.shapes.add_textbox(Inches(0.4), Inches(top_inch), Inches(9.2), Inches(6.5 - top_inch))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(12)
        if index == 0:
            paragraph.font.bold = True


def _patch_existing_slides(prs: Presentation) -> None:
    # 幻灯片 3：项目目标（摘要级）
    slide_goal = prs.slides[2]
    _add_callout_to_slide(
        slide_goal,
        4.2,
        [
            "【A0 本地计算层】DataSource → 生物标志物矩阵 → 序贯事件模型+MCMC 分期 → 文献校准贝叶斯后验",
            "注入 A3/A6/A8/A15；全程本地运行。详见后续「A0 序贯分期与贝叶斯」专节。",
        ],
    )

    # 幻灯片 5：主要创新 — 增补第 4 点
    slide_innov = prs.slides[4]
    _add_callout_to_slide(
        slide_innov,
        5.0,
        [
            "4. A0 本地序贯分期 + 贝叶斯风险融合（新增）",
            "序贯=标志物按顺序依次异常；亚型=恶化顺序不同；共同分期1–5（MCMC）。",
            "贝叶斯：文献先验 P(H|stage) + 队列 z-score 似然比 → 后验，驱动 GP 会诊阈值。",
        ],
    )

    # 幻灯片 8：协同算法技术方案
    slide_tech = prs.slides[7]
    _add_callout_to_slide(
        slide_tech,
        4.5,
        [
            "【A0 增强】在 11 类图谱工具外，新增 tool_疾病分期洞察、tool_贝叶斯风险后验；",
            "与返院 7 阶段流程衔接：阶段/后验概率 → 事件升级与任务分配。",
        ],
    )

    # 幻灯片 9：BPSD 场景
    slide_scene = prs.slides[8]
    _add_callout_to_slide(
        slide_scene,
        4.8,
        [
            "【量化决策】本地序贯分期（NPI/MMSE/MoCA/Barthel）+ BPSD 30天升级后验；",
            "先验来自 JAMA/NeDEM 等文献分层，后验融合返院最新评估 z-score。",
        ],
    )


def _insert_new_section_slides(prs: Presentation) -> int:
    """在末尾追加专节幻灯片，再整体移到技术方案页之后。返回新增页数。"""
    before_count = len(prs.slides)

    _add_textbox_slide(
        prs,
        "A0 本地序贯分期与贝叶斯 — 专节导览",
        [
            "本专节说明算法原理（非单病例报告）",
            "① 序贯含义：多标志物按顺序异常，共用分期轴",
            "② MCMC 推断亚型与阶段 1–5",
            "③ 疾病进展曲线与配图",
            "④ 文献校准先验 + 贝叶斯后验公式",
            "运行：python task-agent/run_staging_pipeline.py",
        ],
        font_size=17,
    )

    _add_textbox_slide(
        prs,
        "序贯分期算法说明",
        [
            "标志物：NPI、MMSE/MoCA 认知逆指标、Barthel/ADL（不含血压/血氧）",
            "每个标志物 = 一个序贯事件：baseline → abnormal，logistic onset",
            "亚型（3种）：事件顺序不同，如 认知→MoCA→NPI→ADL",
            "整体阶段：全体标志物在同一条分期轴 1–5 上，Metropolis-Hastings MCMC 推断",
            "输出：a0_序贯疾病分期结果.json、进展曲线 PNG",
        ],
    )

    _add_textbox_slide(
        prs,
        "核心公式（序贯 + 贝叶斯）",
        [
            "式1  E[L_k|s] = b_k + σ((s−τ_k)/w_k) · (a_k − b_k)   （logistic 期望水平）",
            "式2  log P(D|s) ∝ −½ Σ_k ((x_k − E_k)/σ_k)²   （MCMC 似然）",
            "式3  P(H|stage) ≈ p_s × λ_s   （文献先验，见下页）",
            "式4  odds(H) = P(H)/(1−P(H))",
            "式5  odds(H|D) = odds(H) × ∏_i LR(z_i)",
            "式6  P(H|D) = odds(H|D) / (1 + odds(H|D))",
            "H：30天内 NPI 升高≥4 或任一分项≥4（临床显著）",
        ],
        font_size=15,
    )

    prior_lines = [
        PRIOR_ASSUMPTION_H,
        PRIOR_FORMULA_NOTE,
        "",
        "各阶段文献校准先验 P(H|stage)：",
    ]
    for stage_index in range(1, 6):
        prior_lines.append(
            f"  阶段{stage_index} {STAGE_LABELS[stage_index]}："
            f"{STAGE_BPSD_PRIORS_LITERATURE[stage_index]:.0%} — {STAGE_PRIOR_DERIVATION[stage_index][:55]}…"
        )
    prior_lines.append("")
    prior_lines.append("主要参考文献：")
    for reference in PRIOR_REFERENCES[:4]:
        prior_lines.append(f"[{reference['id']}] {reference['citation'][:70]}…")
    prior_lines.append("[5][6] Aalten 2005；中国 AD 痴呆诊疗指南 2020")

    _add_textbox_slide(prs, "先验概率来源（文献校准）", prior_lines, font_size=13)

    _add_textbox_slide(
        prs,
        "贝叶斯融合与超级 GP 决策阈值",
        [
            "先验：仅由 MCMC 疾病阶段查文献表（如阶段4 → 36%）",
            "似然：各标志物队列 z-score 连续映射为 LR（优于队列↓，差于队列↑）",
            "后验：更新 odds 后与阈值比较",
            "GP 会诊 ≥35%；护士周评 NPI ≥25%；照护员加强监测 ≥20%",
            "工具：tool_贝叶斯风险后验 → A3 会诊 / A8 升级 / A15 报告",
        ],
    )

    for title, filename in FIGURES:
        _add_picture_slide(prs, title, TASK_OUTPUT / filename)

    new_count = len(prs.slides) - before_count

    # 「谢谢」页保持最后：新增页插在谢谢之前
    thanks_index = _find_thanks_slide_index(prs)
    if thanks_index is not None and thanks_index < len(prs.slides) - new_count:
        move_slide(prs, thanks_index, len(prs.slides) - 1)

    return new_count


def main() -> None:
    if not BACKUP_PPT.is_file():
        raise FileNotFoundError(f"未找到原版备份: {BACKUP_PPT}")

    shutil.copy2(BACKUP_PPT, OUTPUT_PPT)
    prs = Presentation(str(OUTPUT_PPT))
    original_count = len(prs.slides)

    _patch_existing_slides(prs)
    added = _insert_new_section_slides(prs)

    prs.save(str(OUTPUT_PPT))

    deliverable = COMPETITION / "交付物" / OUTPUT_PPT.name
    deliverable.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(OUTPUT_PPT, deliverable)

    print(f"源（未改动）: {BACKUP_PPT}")
    print(f"已另存为: {OUTPUT_PPT}")
    print(f"已同步: {deliverable}")
    print(f"幻灯片: {original_count} → {len(prs.slides)}（新增专节 {added} 页，并嵌入摘要/创新/技术/场景页）")


if __name__ == "__main__":
    main()
