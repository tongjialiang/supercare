#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""将 A0 计算方法、核心思想及配图写入技术报告 docx 与比赛 PPT。"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

PROJECT_ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = PROJECT_ROOT / "output"
COMPETITION = Path("/srv/supercare/比赛文档")
PPT_PATH = COMPETITION / "2_PPT_超级[GP-护士-照护员]智能体协同算法.pptx"
REPORT_PATH = COMPETITION / "1_技术报告_超级[GP-护士-照护员]智能体协同算法.docx"

METHOD_SECTION_MCMC = "A0 序贯事件异常模型与 MCMC 疾病进展曲线"

METHOD_PARAGRAPHS_MCMC = [
    "【核心思想—序贯事件模型】每个 BPSD 生物标志物（NPI、MMSE/MoCA 认知逆指标、Barthel/ADL）对应一个序贯事件异常模型："
    "baseline→abnormal，以 logistic onset 刻画该标志物随共同疾病分期 1–5 的期望水平变化；不含收缩压/血氧。",
    "【亚型与共同分期】亚型=标志物异常出现顺序（3 种预设轨迹）；全体标志物处在同一条分期轴上，"
    "由 Metropolis-Hastings MCMC 在队列最新截面上联合推断分期与亚型（非逐指标打分再平均）。",
    "【疾病进展曲线】横轴=疾病分期，纵轴=标志物原始水平；n 个特征 n 条拟合曲线 + 队列分桶经验均值 + 焦点病例标星。",
    "【贝叶斯】先验 P(H) 由 MCMC 分期查表；各标志物 z-score 似然比 odds 更新后验，联动 GP 阈值 35%/25%/20%。",
    "【本地运行】run_staging_pipeline.py + generate_a0_visualizations.py，不调用外部分期 API。",
]

FORMULA_LINES_MCMC = [
    "E[level|stage] = baseline + σ_logistic(stage - onset) × (abnormal - baseline)",
    "log P(obs|stage,subtype) ∝ -Σ (x_i - E_i)² / (2σ_i²)",
    "odds(H) = P(H) / (1-P(H))；odds(H|D) = odds(H) × ∏LR_i",
    "P(H|D) = odds(H|D) / (1 + odds(H|D))",
]

FIGURES = [
    (
        "图A0-0 疾病进展曲线（横轴=分期，纵轴=标志物水平；实线=logistic 拟合，★=焦点病例）",
        "a0_疾病进展曲线_陈女士.png",
    ),
    (
        "图A0-1 序贯进展热图（纵轴=事件触发，横轴=时间，末行=整体分期）",
        "a0_生物标志物分期进展热图_陈女士.png",
    ),
    (
        "图A0-2 序贯事件槽位与 MCMC 整体分期",
        "a0_序贯事件进展图_陈女士.png",
    ),
    (
        "图A0-3 贝叶斯先验与后验（BPSD 30 天升级概率）",
        "a0_贝叶斯先验后验图_陈女士.png",
    ),
]

PPT_SLIDE_MARKERS = [
    "序贯事件模型与 MCMC",
    "图A0-0 疾病进展曲线",
]


def _pick_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "空白" in name:
            return layout
    return prs.slide_layouts[0]


def _ppt_already_has_mcmc_section(prs: Presentation) -> bool:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            if any(marker in text for marker in PPT_SLIDE_MARKERS):
                return True
    return False


def update_technical_report() -> None:
    doc = Document(str(REPORT_PATH))
    if any(METHOD_SECTION_MCMC in (paragraph.text or "") for paragraph in doc.paragraphs):
        print("技术报告已含 MCMC 进展曲线章节，跳过重复写入")
        return

    doc.add_page_break()
    doc.add_heading(METHOD_SECTION_MCMC, level=1)
    doc.add_paragraph(
        "分期范式：每个生物标志物一个序贯事件异常模型 + MCMC 推断共同分期；"
        "主可视化为一组疾病进展曲线（横轴分期、纵轴水平）。"
    )

    for paragraph_text in METHOD_PARAGRAPHS_MCMC:
        paragraph = doc.add_paragraph(paragraph_text)
        paragraph.paragraph_format.space_after = Pt(8)

    doc.add_heading("计算公式", level=2)
    for line in FORMULA_LINES_MCMC:
        doc.add_paragraph(f"• {line}")

    doc.add_heading("配图", level=2)
    for caption, filename in FIGURES:
        image_path = OUTPUT_DIR / filename
        doc.add_paragraph(caption)
        if image_path.is_file():
            doc.add_picture(str(image_path), width=Inches(5.8))
        else:
            doc.add_paragraph(f"（未找到：{filename}）")

    doc.save(str(REPORT_PATH))
    print(f"技术报告已更新: {REPORT_PATH}")


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layouts = list(prs.slide_layouts)
    layout = layouts[1] if len(layouts) > 1 else layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body is None and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
    if body:
        text_frame = body.text_frame
        text_frame.clear()
        for index, line in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.font.size = PptPt(16)


def _add_picture_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(_pick_blank_layout(prs))
    textbox = slide.shapes.add_textbox(PptInches(0.3), PptInches(0.15), PptInches(9.4), PptInches(0.7))
    textbox.text_frame.text = title
    textbox.text_frame.paragraphs[0].font.size = PptPt(22)
    textbox.text_frame.paragraphs[0].font.bold = True
    if image_path.is_file():
        slide.shapes.add_picture(str(image_path), PptInches(0.25), PptInches(0.95), width=PptInches(9.2))


def update_ppt() -> None:
    prs = Presentation(str(PPT_PATH))
    if _ppt_already_has_mcmc_section(prs):
        print("PPT 已含 MCMC/进展曲线页，跳过重复追加")
        return

    before = len(prs.slides)

    _add_bullet_slide(
        prs,
        "序贯事件模型与 MCMC（本地实现）",
        [
            "每标志物：baseline→abnormal + logistic onset，横轴=共同分期 1–5",
            "Metropolis-Hastings MCMC 推断亚型与整体分期",
            "主图：疾病进展曲线（n 特征 n 条，纵轴=原始水平）",
            "BPSD 指标：NPI / MMSE / MoCA / Barthel，不含血压血氧",
            "贝叶斯：分期先验 + 标志物 LR → 后验，联动 GP 阈值",
        ],
    )
    for caption, filename in FIGURES:
        short_title = caption.split("（")[0]
        _add_picture_slide(prs, short_title, OUTPUT_DIR / filename)

    prs.save(str(PPT_PATH))
    print(f"PPT 已更新: {PPT_PATH}（{before} → {len(prs.slides)} 页）")


def main() -> None:
    update_technical_report()
    update_ppt()


if __name__ == "__main__":
    main()
